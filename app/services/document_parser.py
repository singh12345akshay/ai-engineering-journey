import pymupdf
from pptx import Presentation
from typing import List
import os


def extract_text_from_pdf(file_path: str) -> List[dict]:
    """
    Extract text from PDF file page by page.
    Returns list of dicts with text and page metadata.
    """
    pages = []

    doc = pymupdf.open(file_path)

    for page_num, page in enumerate(doc):
        text = page.get_text()

        # skip empty pages
        if not text.strip():
            continue

        pages.append({
            "text": text.strip(),
            "page_number": page_num + 1,
            "source": os.path.basename(file_path),
            "type": "pdf"
        })

    doc.close()
    return pages


def extract_text_from_pptx(file_path: str) -> List[dict]:
    """
    Extract text from PPTX file slide by slide.
    Returns list of dicts with text and slide metadata.
    """
    slides = []

    prs = Presentation(file_path)

    for slide_num, slide in enumerate(prs.slides):
        slide_text = []

        for shape in slide.shapes:
            # only process shapes that have text
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        full_text = "\n".join(slide_text)

        # skip empty slides
        if not full_text.strip():
            continue

        slides.append({
            "text": full_text,
            "slide_number": slide_num + 1,
            "source": os.path.basename(file_path),
            "type": "pptx"
        })

    return slides


def chunk_text(text: str, chunk_size: int = 500,
               overlap: int = 50) -> List[str]:
    """
    Split text into overlapping chunks by word count.

    chunk_size: target words per chunk
    overlap: words shared between consecutive chunks

    Why word count not character count:
    Word boundaries are more natural split points.
    500 words ≈ 375 tokens ≈ fits well in context window.
    """
    words = text.split()
    chunks = []

    if len(words) <= chunk_size:
        # text is short enough — no chunking needed
        return [text]

    start = 0
    while start < len(words):
        end = start + chunk_size
        chunk = " ".join(words[start:end])
        chunks.append(chunk)

        # move forward by chunk_size minus overlap
        # this creates the overlap between chunks
        start += chunk_size - overlap

        # stop if remaining text is too small
        if start >= len(words):
            break

    return chunks


def parse_and_chunk_document(file_path: str,
                              chunk_size: int = 500,
                              overlap: int = 50,
                              original_filename: str = None) -> List[dict]:
    """
    Parse document and return chunked pieces ready for embedding.
    original_filename: the real filename to store in metadata
                       instead of the temp file path
    """
    ext = os.path.splitext(file_path)[1].lower()
    display_name = original_filename or os.path.basename(file_path)

    if ext == ".pdf":
        pages = extract_text_from_pdf(file_path)
    elif ext in [".pptx", ".ppt"]:
        pages = extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    # override source name with real filename
    for page in pages:
        page["source"] = display_name

    chunks = []
    chunk_index = 0

    for page in pages:
        page_chunks = chunk_text(page["text"], chunk_size, overlap)

        for chunk in page_chunks:
            if len(chunk.split()) < 10:
                continue

            chunk_data = {
                "text": chunk,
                "metadata": {
                    "source": page["source"],
                    "type": page["type"],
                    "chunk_index": chunk_index,
                    "page_or_slide": page.get("page_number") or page.get("slide_number") or page.get("sheet_number", 0)
                }
            }
            chunks.append(chunk_data)
            chunk_index += 1

    return chunks